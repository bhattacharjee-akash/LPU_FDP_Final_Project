import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define color scheme (LPU Orange & Charcoal theme)
    ORANGE = RGBColor(227, 108, 10)  # LPU Orange
    CHARCOAL = RGBColor(38, 38, 38)  # Charcoal Dark
    LIGHT_GRAY = RGBColor(245, 245, 245)
    
    slides_data = [
        # Slide 1: Title
        {
            "layout": 0,
            "title": "LPU Academic Copilot",
            "subtitle": "Multi-Agent AI Platform for Faculty Workflow Automation\nLovely Professional University (LPU) FDP Capstone",
            "notes": "Good morning members of the FDP evaluation panel. Today, we present the LPU Academic Copilot, a full-stack platform driven by a 10-agent orchestrator to automate the standard curriculum preparation workflow."
        },
        # Slide 2: Problem Statement
        {
            "layout": 1,
            "title": "Problem Statement",
            "bullets": [
                "Faculty spend significant hours manually drafting weekly delivery schedules.",
                "Designing assignments and quiz banks with correct answers is tedious and error-prone.",
                "Reviewing exam questions manually for cognitive weight (Bloom's Taxonomy) is highly subjective.",
                "Verifying syllabus compliance takes multiple iteration reviews, causing process delays."
            ],
            "notes": "The primary problem is administrative overload. Faculty members spend over 8 hours per planning cycle on paperwork, taking away from direct student mentoring."
        },
        # Slide 3: Motivation
        {
            "layout": 1,
            "title": "Motivation",
            "bullets": [
                "Empower lovely faculty members by automating repetitive documentation.",
                "Maintain absolute standardization across academic packages.",
                "Provide objective, automated feedback on curriculum compliance prior to class delivery.",
                "Accelerate institutional agility in resource compilation."
            ],
            "notes": "We wanted to leverage modern agentic workflows to build an automated assistant that serves as a quality-audited companion to faculty members."
        },
        # Slide 4: Objectives
        {
            "layout": 1,
            "title": "Objectives",
            "bullets": [
                "Parse standard syllabus PDF files dynamically.",
                "Orchestrate a 10-Agent pipeline in a concurrent/dependent sequence.",
                "Generate structured lecture schedules, progressive assignments, and 20 MCQs.",
                "Conduct objective Bloom's cognitive mapping and Course Outcome correlation.",
                "Output a publication-ready Course Package PDF containing all reviewed deliverables."
            ],
            "notes": "Our objectives focus on automation, structuring, mapping, and compilation, delivering a complete packet in under 40 seconds."
        },
        # Slide 5: Existing Challenges
        {
            "layout": 1,
            "title": "Existing Challenges in AI Document Synthesis",
            "bullets": [
                "Invalid Model Names: Out-of-date model identifiers cause live API crashes.",
                "DNS & Network Barriers: Direct PostgreSQL IPv6 calls fail inside Render's IPv4 VMs.",
                "Query Param Incompatibilities: Unused parameters (like pgbouncer) crash DB drivers.",
                "Download Authentication Conflicts: Browser file downloads lack JWT bearer headers."
            ],
            "notes": "During deployment, we solved major real-world challenges, such as handling Render's IPv4 network limitations when connecting to Supabase's IPv6 databases, and fixing download authentication loops."
        },
        # Slide 6: Proposed Solution
        {
            "layout": 1,
            "title": "Proposed Solution",
            "bullets": [
                "Full-stack web application linking Next.js 15 and FastAPI.",
                "State-of-the-art Google Gemini 1.5 Flash API reasoning models.",
                "An asynchronous 10-node agent orchestrator running in background threads.",
                "Supabase PostgreSQL and Object Storage database backend.",
                "CORS-compliant and authenticated APIs with public token-less download streams."
            ],
            "notes": "Our solution decoupling Next.js and FastAPI provides a seamless UI tracking timeline while running intensive multi-agent models in the background."
        },
        # Slide 7: System Architecture
        {
            "layout": 1,
            "title": "System Architecture & Data Flows",
            "bullets": [
                "Frontend: Next.js 15, TypeScript, Tailwind CSS, and Supabase client-side Auth.",
                "Backend: FastAPI (Python 3.12/3.14) serving REST APIs and managing Alembic migrations.",
                "Database Connection: Powered by Supabase Connection Pooler (Port 6543) resolving to IPv4.",
                "Authentication: Supabase JWT validation inside FastAPI dependencies."
            ],
            "notes": "The system architecture is built for resiliency. We use the Supabase Connection Pooler over port 6543 to bypass Render network restrictions."
        },
        # Slide 8: Multi-Agent Workflow
        {
            "layout": 1,
            "title": "The 10-Agent Orchestrator Pipeline",
            "bullets": [
                "Planning Node: Extracts syllabus structure, units, and Course Outcomes.",
                "Generation Group (Parallel): Lesson Plan, Assignment, Quiz, and Question Paper agents.",
                "Mapping Group: Bloom Taxonomy Evaluator and Course Outcomes (CO) alignment nodes.",
                "Review Group: Consistency checker (Reviewer) and Compliance Auditor (Quality Agent).",
                "Compilation Node: Compiles raw structures into a styled ReportLab PDF pack."
            ],
            "notes": "Each agent is specialized. By prompting them in a sequence, the output of the Planning Agent forms the context for the generation, mapping, and quality check agents."
        },
        # Slide 9: Technology Stack
        {
            "layout": 1,
            "title": "Technology Stack",
            "bullets": [
                "Frontend: React 19, Next.js 15 (App Router), Framer Motion, Lucide Icons",
                "Backend: FastAPI, Uvicorn ASGI Server, SQLAlchemy ORM",
                "Database & Storage: Supabase PostgreSQL (PgBouncer), Supabase Storage Buckets",
                "AI Models: Google Gemini 1.5 Flash (Primary), Groq Cloud APIs (Secondary)",
                "PDF Engine: ReportLab layout flowables"
            ],
            "notes": "We are using modern, fast frameworks. Google Gemini 1.5 Flash provides high reasoning speeds, and ReportLab enables dynamic PDF layouts."
        },
        # Slide 10: Database Design
        {
            "layout": 1,
            "title": "Relational Database Design",
            "bullets": [
                "users: Stores synchronized authentication profiles.",
                "syllabi: Holds the parsed text, code, and PDF path parameters.",
                "generation_histories: Tracks processing and completion states.",
                "agent_execution_logs: Feeds the live timeline in the frontend.",
                "All tables use Cascade On Delete constraints for data integrity."
            ],
            "notes": "The database contains six main relational tables. The agent logs feed our front-end progress indicators, providing trace information to users in real-time."
        },
        # Slide 11: Key Features
        {
            "layout": 1,
            "title": "Key Platform Features",
            "bullets": [
                "Interactive Uploader: Simple drag & drop interface for PDF files.",
                "Live Timeline Logs: Real-time progress updates directly from background workers.",
                "Calibration Settings: Fine-tune model parameters and faculty details.",
                "Public Downloads: Direct, secure downloads without header conflicts.",
                "Resilient DB Fallback: Gracefully routes to SQLite locally if remote PG is offline."
            ],
            "notes": "The key features include progress tracking, configuration customization, and automatic database fallbacks, rendering the platform robust."
        },
        # Slide 12: Application Demonstration Flow
        {
            "layout": 1,
            "title": "Application Demonstration Flow",
            "bullets": [
                "Step 1: User signs up/in on the Vercel frontend portal.",
                "Step 2: Uploads a standard syllabus PDF (e.g. Machine Learning).",
                "Step 3: Trigger Multi-Agent Workflow triggers background tasks.",
                "Step 4: Status timelines trace agent outputs concurrently.",
                "Step 5: The final Course Pack PDF is downloaded with 1 click."
            ],
            "notes": "For the evaluation, we demonstrate a complete flow using a sample Machine Learning syllabus. We register, upload, watch the log timeline, and download the output PDF."
        },
        # Slide 13: Results and Benefits
        {
            "layout": 1,
            "title": "Results & Benefits",
            "bullets": [
                "Over 8 hours saved per course planning cycle.",
                "Ensures 100% standard formatting across university course packs.",
                "Objective grading: Bloom's Taxonomy cognitive weights are mapped uniformly.",
                "Immediate quality feedback allows faster curriculum iterations."
            ],
            "notes": "Our metrics show a drastic improvement in efficiency, with 100% adherence to standard formatting requirements."
        },
        # Slide 14: Future Scope
        {
            "layout": 1,
            "title": "Future Scope & Enhancements",
            "bullets": [
                "Direct LPU LMS / UMS Integration: Push course packs to class groups instantly.",
                "AI Lecture Notes: Automatically generate slide decks and lecture notes for weekly plans.",
                "Syllabus Comparison: Cross-correlate multiple syllabi to identify gaps.",
                "OpenAI/Anthropic Engine integrations."
            ],
            "notes": "In the future, we plan to connect this directly to LPU's LMS and automate lecture slide synthesis."
        },
        # Slide 15: Security & Performance
        {
            "layout": 1,
            "title": "Security & Performance Considerations",
            "bullets": [
                "Data Isolation: Users can only see and manage their own syllabus histories.",
                "CORS Protection: Domain locks prevent cross-origin scripting hacks.",
                "Async Background Workers: Keeps the API gateway responsive during heavy LLM calls.",
                "Sanitized Database URLs: Dynamic filters prevent credential leaks in trace logs."
            ],
            "notes": "Security is maintained via RLS policies and CORS domain restrictions, while asynchronous execution keeps the app highly responsive."
        },
        # Slide 16: Conclusion
        {
            "layout": 1,
            "title": "Conclusion",
            "bullets": [
                "Decoupled multi-agent architecture successfully validated in production.",
                "Delivers a complete, audited course package dynamically.",
                "Bypasses complex deployment and database connection conflicts.",
                "Innovative solution to elevate institutional planning efficiency."
            ],
            "notes": "In conclusion, the LPU Academic Copilot shows how agentic workflows can transform administrative tasks into standard, audited outputs."
        },
        # Slide 17: Thank You
        {
            "layout": 0,
            "title": "Thank You",
            "subtitle": "Lovely Professional University (LPU)\nFDP Final Evaluation - Q&A Session",
            "notes": "Thank you for your time. I am open to any questions you may have."
        }
    ]
    
    # Process slides
    for idx, slide_info in enumerate(slides_data):
        layout_idx = slide_info["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Add presenter notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_info.get("notes", "")
        
        # Add Title
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]
        
        # Add subtitle (for layout 0) or bullets (for layout 1)
        if layout_idx == 0:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = slide_info["subtitle"]
        elif layout_idx == 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            for b_idx, bullet in enumerate(slide_info.get("bullets", [])):
                p = tf.add_paragraph() if b_idx > 0 else tf.paragraphs[0]
                p.text = bullet
                p.level = 0
                
    # Save presentation
    output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "LPU_FDP_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
